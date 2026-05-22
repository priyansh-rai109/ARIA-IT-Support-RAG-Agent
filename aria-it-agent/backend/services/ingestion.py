import os
import queue
import threading
import traceback
import csv
from datetime import datetime
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma

from backend.config import SOURCES_DIR, CHROMA_DB_DIR, EMBEDDING_MODEL_NAME
from backend.database.connection import get_db_connection

# Ingestion queue and thread variables
ingestion_queue = queue.Queue()
db_lock = threading.Lock()
embeddings_model = None

def get_embeddings():
    global embeddings_model
    if embeddings_model is None:
        embeddings_model = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL_NAME)
    return embeddings_model

def get_vectorstore():
    return Chroma(persist_directory=CHROMA_DB_DIR, embedding_function=get_embeddings())

# Helper to update status in DB
def update_status(file_name, status, progress, error_message=None):
    with db_lock:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO ingestion_status (file_name, status, progress, error_message, timestamp)
            VALUES (?, ?, ?, ?, ?)
            ON CONFLICT(file_name) DO UPDATE SET
                status=excluded.status,
                progress=excluded.progress,
                error_message=excluded.error_message,
                timestamp=excluded.timestamp
        """, (file_name, status, progress, error_message, datetime.now().isoformat()))
        conn.commit()
        conn.close()

# Parsers
def parse_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def parse_pdf(file_path):
    reader = PdfReader(file_path)
    text = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            text.append(t)
    return "\n\n".join(text)

def parse_docx(file_path):
    doc = DocxDocument(file_path)
    text = []
    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text)
    # Extract tables if present
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                text.append(" | ".join(row_text))
    return "\n\n".join(text)

def parse_pptx(file_path):
    prs = PptxPresentation(file_path)
    text = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_text = [f"--- Slide {slide_idx+1} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        text.append("\n".join(slide_text))
    return "\n\n".join(text)

def parse_csv(file_path):
    text = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        headers = next(reader, None)
        if headers:
            text.append(f"Headers: {', '.join(headers)}")
        for idx, row in enumerate(reader):
            text.append(f"Row {idx+1}: {', '.join(row)}")
    return "\n".join(text)

def extract_file_content(file_path):
    _, ext = os.path.splitext(file_path.lower())
    if ext == ".txt":
        return parse_txt(file_path)
    elif ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".csv":
        return parse_csv(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")

# worker thread processing ingestion
def process_ingestion_queue():
    while True:
        try:
            file_name, file_path = ingestion_queue.get()
            if file_name is None: # poison pill
                break
            
            update_status(file_name, "Processing", 20)
            
            try:
                # 1. Extract content
                text = extract_file_content(file_path)
                update_status(file_name, "Processing", 50)
                
                # 2. Dynamic Chunking
                # Recursive Character Splitter handles structural trimmers (paragraphs -> sentences -> chars)
                splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
                doc = Document(page_content=text, metadata={"id": file_name, "title": file_name, "category": "Ingested Document"})
                chunks = splitter.split_documents([doc])
                update_status(file_name, "Processing", 70)
                
                # 3. Vectorize and Ingest to VectorDB
                vectorstore = get_vectorstore()
                vectorstore.add_documents(chunks)
                
                # 4. Notify RAG pipeline of index update (triggers BM25 recreation if active)
                from backend.services.rag_pipeline import notify_db_update
                notify_db_update()
                
                update_status(file_name, "Complete", 100)
                
            except Exception as e:
                err_msg = traceback.format_exc()
                print(f"Error ingesting {file_name}: {err_msg}")
                update_status(file_name, "Failed", 0, str(e))
                
            finally:
                ingestion_queue.task_done()
        except Exception as queue_err:
            print(f"Queue processor encountered severe error: {queue_err}")

# API to enqueue document for ingestion
def enqueue_document(file_name, file_bytes):
    os.makedirs(SOURCES_DIR, exist_ok=True)
    file_path = os.path.join(SOURCES_DIR, file_name)
    
    # Save the file locally
    with open(file_path, "wb") as f:
        f.write(file_bytes)
        
    update_status(file_name, "Pending", 0)
    ingestion_queue.put((file_name, file_path))
    return {"message": f"Successfully queued {file_name} for ingestion.", "status": "Pending"}

# Start background worker thread
worker_thread = threading.Thread(target=process_ingestion_queue, daemon=True)
worker_thread.start()
